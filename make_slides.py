#!/usr/bin/env python3
"""Generate AutoFocus project slides as an editable .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------- colour palette ----------
BG      = RGBColor(0x0D, 0x11, 0x17)
BLUE    = RGBColor(0x58, 0xA6, 0xFF)
ORANGE  = RGBColor(0xFF, 0x8C, 0x00)
WHITE   = RGBColor(0xC9, 0xD1, 0xD9)
GREY    = RGBColor(0x8B, 0x94, 0x9E)
DARKBOX = RGBColor(0x16, 0x1B, 0x22)

BAND_COLORS = {
    'u': RGBColor(0x81, 0x72, 0xB3),
    'g': RGBColor(0x46, 0x82, 0xB4),
    'r': RGBColor(0xFF, 0x8C, 0x00),
    'i': RGBColor(0xC4, 0x4E, 0x52),
    'z': RGBColor(0x93, 0x78, 0x60),
}

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ---- helpers ----
def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def set_text(tf, text, size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def add_para(tf, text, size=22, color=WHITE, bold=False, space_before=6, bullet=False, level=0, alignment=None):
    p = tf.add_paragraph()
    p.level = level
    p.space_before = Pt(space_before)
    if alignment is not None:
        p.alignment = alignment
    if bullet:
        p.bullet = True  # might not work in all readers, we use prefix instead
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)
    tb = add_textbox(slide, 1, 2.0, 11.3, 1.5)
    set_text(tb.text_frame, title, size=44, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        tb2 = add_textbox(slide, 1, 3.6, 11.3, 1.0)
        set_text(tb2.text_frame, subtitle, size=24, color=GREY, alignment=PP_ALIGN.CENTER)
    return slide

def section_slide(heading):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    tb = add_textbox(slide, 0.8, 0.4, 11.7, 0.9)
    set_text(tb.text_frame, heading, size=36, color=BLUE, bold=True)
    return slide

def add_bullets(slide, items, left=1.0, top=1.6, width=11.3, height=5.0, size=22, color=WHITE):
    tb = add_textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            set_text(tf, "•  " + item, size=size, color=color)
        else:
            add_para(tf, "•  " + item, size=size, color=color, space_before=10)
    return tb

def add_image_placeholder(slide, label, left, top, width, height):
    """Add a labelled box where user can insert an image."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARKBOX
    shape.line.color.rgb = GREY
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    set_text(tf, label, size=16, color=GREY, alignment=PP_ALIGN.CENTER)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

# =====================================================================
# SLIDE 1 — Title
# =====================================================================
s = title_slide("🔭 AutoFocus",
                "Real-Time Seeing & Focus Monitor for the Bok 90Prime")
tb = add_textbox(s, 1, 4.8, 11.3, 1.0)
set_text(tb.text_frame, "Jenny  |  Steward Observatory, University of Arizona\nApril 2026",
         size=20, color=GREY, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 2 — Motivation
# =====================================================================
s = section_slide("Why Do We Need This?")
add_bullets(s, [
    "Seeing changes constantly during a night — but you don't know how much until you reduce the data",
    "Focus drifts with temperature; manual focus checks waste precious telescope time",
    "The PI back home has no visibility into real-time data quality",
    "Existing workflow: DS9 + eyeball → slow, subjective, not quantitative",
])
tb = add_textbox(s, 1, 5.6, 11.3, 0.8)
set_text(tb.text_frame,
         "Goal: Automated, quantitative, real-time PSF monitoring — raw FITS → interactive dashboard in < 30 s",
         size=22, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 3 — What It Does
# =====================================================================
s = section_slide("What AutoFocus Does")
add_bullets(s, [
    "📸  Watches for new exposures as they arrive (watchdog)",
    "⚡  Bias-subtracts, flat-fields, runs source extraction (SEP)",
    "⭐  Measures FWHM, ellipticity, and airmass per frame",
    "📊  Generates live interactive Plotly dashboard",
    "🌐  Shares via HTTP + ngrok — accessible from anywhere",
], top=1.6, width=7)

# stat boxes on the right
for i, (val, label) in enumerate([("~10 s", "per frame"), ("5", "bands"), ("8", "amps / CCD")]):
    shape = slide_shapes = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.0), Inches(1.8 + i * 1.7), Inches(3.2), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARKBOX
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    set_text(tf, val, size=36, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_para(tf, label, size=16, color=GREY, alignment=PP_ALIGN.CENTER)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER

# =====================================================================
# SLIDE 4 — Pipeline Architecture
# =====================================================================
s = section_slide("Pipeline Architecture")

steps = ["📁 New FITS\n(watchdog)", "🔧 Calibration\n(bias / flat)", "🔍 SEP Extract\n(× 8 amps)",
         "📈 GMM Select\n(stars vs gal)", "📊 Dashboard\n(Plotly HTML)"]
for i, step in enumerate(steps):
    x = 0.6 + i * 2.55
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(1.8), Inches(2.2), Inches(1.4))
    shape.fill.solid(); shape.fill.fore_color.rgb = DARKBOX
    shape.line.color.rgb = GREY; shape.line.width = Pt(1)
    tf = shape.text_frame; tf.word_wrap = True
    set_text(tf, step, size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
    for p in tf.paragraphs: p.alignment = PP_ALIGN.CENTER
    # arrow
    if i < len(steps) - 1:
        atb = add_textbox(s, x + 2.2, 2.1, 0.4, 0.8)
        set_text(atb.text_frame, "→", size=28, color=ORANGE, alignment=PP_ALIGN.CENTER)

add_bullets(s, [
    "Incremental mode: caches master calibrations & per-file detections — only new frames are reduced",
    "Multi-band: --filter all auto-discovers bands and matches each frame to the correct flat",
    "Star/galaxy separation: 2-component GMM on FWHM, flux ratio, magnitude, ellipticity",
], top=3.8, size=20)

# =====================================================================
# SLIDE 5 — GMM for Star/Galaxy Separation (Deep Dive)
# =====================================================================
s = section_slide("Gaussian Mixture Model — How It Works")

# Left: explanation
add_bullets(s, [
    "Problem: SEP detects everything — stars, galaxies, cosmic rays, hot pixels. "
    "We need only isolated, unsaturated stars to measure the PSF.",
    "Solution: Gaussian Mixture Model (GMM) with K = 2 components. "
    "Each component is a multivariate Gaussian in feature space.",
    "Features used (per source): FWHM, flux_radius / FWHM (concentration), "
    "instrumental magnitude, ellipticity → 4-D feature vector.",
    "GMM learns two clusters — the tight, round, compact cluster is \"stars\"; "
    "the diffuse, extended cluster is \"everything else.\"",
    "Each source gets a posterior probability P(star | features). "
    "We keep sources with P > 0.5 for PSF statistics.",
    "Why GMM over a hard cut? Stars and galaxies overlap in any single parameter. "
    "GMM finds the optimal boundary in multi-dimensional space — "
    "robust even when seeing varies frame-to-frame.",
], top=1.5, size=19)

# =====================================================================
# SLIDE 6 — GMM: Mathematical Formulation
# =====================================================================
s = section_slide("GMM — Mathematical Formulation")

tb = add_textbox(s, 0.8, 1.5, 11.7, 5.5)
tf = tb.text_frame; tf.word_wrap = True
set_text(tf,
    "The probability density of a K-component GMM:", size=20, color=WHITE)
add_para(tf, "", size=10, color=WHITE)
add_para(tf,
    "    p(x) = Σ_k  π_k · N(x | μ_k, Σ_k)      k = 1, ..., K", size=22, color=ORANGE, bold=True)
add_para(tf, "", size=10, color=WHITE)
add_para(tf, "where:", size=20, color=WHITE)
add_para(tf, "    π_k  = mixing weight (fraction of sources in component k),  Σ π_k = 1", size=18, color=WHITE)
add_para(tf, "    μ_k  = mean vector (centroid in feature space)", size=18, color=WHITE)
add_para(tf, "    Σ_k  = covariance matrix (shape & orientation of the cluster)", size=18, color=WHITE)
add_para(tf, "", size=10, color=WHITE)
add_para(tf, "Parameters are fit via Expectation-Maximization (EM):", size=20, color=WHITE)
add_para(tf, "    E-step: compute responsibility  r_k(x_i) = π_k N(x_i|μ_k,Σ_k) / p(x_i)", size=18, color=ORANGE)
add_para(tf, "    M-step: update  π_k, μ_k, Σ_k  from weighted statistics", size=18, color=ORANGE)
add_para(tf, "", size=10, color=WHITE)
add_para(tf,
    "In AutoFocus: K = 2 (stars vs non-stars), full covariance, "
    "scikit-learn GaussianMixture, converges in < 0.1 s per frame.",
    size=18, color=GREY)

# =====================================================================
# SLIDE 7 — Why GMM Works for Star Selection
# =====================================================================
s = section_slide("Why GMM Works for Star Selection")

add_bullets(s, [
    "Stars are point sources convolved with the PSF → they form a TIGHT cluster: "
    "narrow FWHM range, low ellipticity, high concentration (flux_radius ≈ FWHM).",
    "Galaxies are extended → BROADER distribution: larger FWHM, lower concentration, "
    "higher ellipticity. Cosmic rays are too sharp (FWHM ≪ seeing).",
    "In 4-D feature space, these populations are well-separated Gaussians "
    "→ GMM is a natural, principled model for this bimodal distribution.",
    "Advantages over hard cuts (e.g., FWHM < 5 & e < 0.3):",
    "     • Adapts automatically to changing seeing — no manual threshold tuning",
    "     • Accounts for correlations between features (via full covariance)",
    "     • Provides soft probabilities, not binary labels — can tune purity vs completeness",
    "     • Works consistently across bands with different PSF sizes",
], top=1.5, size=19)

# =====================================================================
# SLIDE 8 — Application to JWST PSF Studies
# =====================================================================
s = section_slide("GMM for JWST PSF Characterization")

add_bullets(s, [
    "JWST's PSF is complex: diffraction-limited core + hexagonal diffraction spikes "
    "+ wavelength-dependent Airy rings. It varies across the field (OPD maps) and with time "
    "(thermal breathing of the primary mirror segments).",
    "Challenge: JWST resolves galaxies even at z > 2 — the star/galaxy overlap "
    "is MUCH worse than ground-based. Simple morphological cuts fail.",
    "GMM in JWST context:",
    "     • Features: FWHM, concentration index (C = r80/r20), Sérsic index, ellipticity, "
    "       encircled energy fraction, sharpness, roundness",
    "     • Stars form a tight locus in this high-D space; compact galaxies scatter broadly",
    "     • Works per-detector, per-filter, per-epoch — captures PSF spatial & temporal variation",
], top=1.5, size=19)

tb = add_textbox(s, 0.8, 5.5, 11.7, 1.5)
tf = tb.text_frame; tf.word_wrap = True
set_text(tf, "Why this matters for JWST science:", size=20, color=BLUE, bold=True)
add_para(tf, "•  Weak lensing (COSMOS-Web, PRIMER): PSF model errors dominate shear systematics. "
             "Clean star samples are essential for PSF interpolation (PSFEx, Piff, WebbPSF).", size=17, color=WHITE)
add_para(tf, "•  Crowded fields (globular clusters, LMC): thousands of stars per chip — "
             "GMM scales well and avoids manual tuning per pointing.", size=17, color=WHITE)
add_para(tf, "•  Time-domain (transients, moving objects): PSF changes between visits — "
             "GMM re-fits automatically each epoch with no prior assumptions.", size=17, color=WHITE)

# =====================================================================
# SLIDE 9 — GMM: Ground vs Space Comparison
# =====================================================================
s = section_slide("GMM Star Selection: Ground vs JWST")

comp_data = [
    ("", "Ground-Based (Bok 90Prime)", "Space (JWST NIRCam)"),
    ("PSF shape", "Seeing-dominated Gaussian\n(FWHM ~ 1–2\")", "Diffraction-limited Airy\n(FWHM ~ 0.03–0.16\")"),
    ("Star/galaxy\noverlap", "Moderate — galaxies\nclearly extended", "Severe — compact galaxies\nmimic point sources"),
    ("Features", "FWHM, concentration,\nmag, ellipticity (4-D)", "FWHM, C, Sérsic, e,\nEE fraction, sharpness (6-D+)"),
    ("PSF variation", "Slowly varies with\nseeing & airmass", "Varies across field (OPD)\n& with time (thermal)"),
    ("GMM advantage", "Adapts to changing\nseeing per frame", "Captures per-detector,\nper-filter PSF locus"),
]
rows, cols = len(comp_data), 3
tbl = s.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.0)).table
tbl.columns[0].width = Inches(2.2)
tbl.columns[1].width = Inches(4.75)
tbl.columns[2].width = Inches(4.75)
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = comp_data[r][c]
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARKBOX if r == 0 else BG
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(16)
                run.font.bold = (r == 0 or c == 0)
                run.font.color.rgb = BLUE if r == 0 else (ORANGE if c == 0 else WHITE)

# =====================================================================
# SLIDE 10 — Key Measurements
# =====================================================================
s = section_slide("Key Measurements")

table_data = [
    ("Metric", "Method", "Purpose"),
    ("FWHM", "Radial profile → Gaussian fit", "Seeing & focus quality"),
    ("Ellipticity", "e = 1 − b/a from SEP moments", "Tracking / wind / optics"),
    ("Airmass", "FITS header keyword", "Seeing vs zenith angle"),
    ("PSF contour", "Median-stacked star cutouts", "Visual PSF shape diagnostic"),
    ("Focus curve", "Parabola fit: FWHM vs LVDTC", "Best focus position"),
]
rows, cols = len(table_data), 3
tbl = s.shapes.add_table(rows, cols, Inches(1), Inches(1.6), Inches(11.3), Inches(4)).table
for c in range(cols):
    tbl.columns[c].width = Inches([2.5, 4.5, 4.3][c])
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARKBOX if r == 0 else BG
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = BLUE if r == 0 else WHITE
                run.font.bold = (r == 0)

tb = add_textbox(s, 1, 6.2, 11.3, 0.5)
set_text(tb.text_frame, "Plate scale: 0.455\"/pixel (Bok 90Prime)  |  FWHM reported in arcsec",
         size=16, color=GREY, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 6 — Multi-Band Support
# =====================================================================
s = section_slide("Multi-Band Support")
tb = add_textbox(s, 1, 1.3, 11, 0.5)
set_text(tb.text_frame, "--filter all → auto-discovers every band in the data", size=20, color=WHITE)

band_table = [
    ("Band", "Frames", "Median FWHM (pix)", "Mean Airmass"),
    ("u", "5", "3.45", "1.26"),
    ("g", "45", "3.69", "1.08"),
    ("r", "61", "3.03", "1.06"),
    ("i", "5", "2.55", "1.26"),
    ("z", "5", "2.78", "1.26"),
]
rows, cols = len(band_table), 4
tbl = s.shapes.add_table(rows, cols, Inches(2), Inches(2.2), Inches(9.3), Inches(3.5)).table
for c in range(cols):
    tbl.columns[c].width = Inches([1.5, 1.8, 3.0, 3.0][c])
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = band_table[r][c]
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARKBOX if r == 0 else BG
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(20)
                run.font.bold = (r == 0)
                if r == 0:
                    run.font.color.rgb = BLUE
                elif c == 0:
                    band = band_table[r][0]
                    run.font.color.rgb = BAND_COLORS.get(band, WHITE)
                    run.font.bold = True
                else:
                    run.font.color.rgb = WHITE

tb = add_textbox(s, 1, 6.0, 11.3, 0.6)
set_text(tb.text_frame,
         "g-band FWHM ~19% larger than r-band → chromatic focus offset from refractive corrector",
         size=18, color=GREY, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 7 — Interactive Dashboard
# =====================================================================
s = section_slide("Interactive Dashboard")
add_bullets(s, [
    "3-panel layout: FWHM · Airmass · Ellipticity vs time",
    "Hover: filename, time, FWHM, filter, airmass",
    "PSF contour popup: median-stacked star image with scale bar",
    "Auto-refresh every 30 s with countdown timer",
    "Color-coded by filter band",
], width=6)
# Placeholder for screenshot
img_path = os.path.join(os.path.dirname(__file__), "focus_output", "focus_time_series.png")
if os.path.exists(img_path):
    s.shapes.add_picture(img_path, Inches(7.2), Inches(1.4), Inches(5.5))
else:
    add_image_placeholder(s, "Insert screenshot:\nfocus_time_series.png", 7.2, 1.4, 5.5, 4.5)

# =====================================================================
# SLIDE 8 — Real-Time Monitoring
# =====================================================================
s = section_slide("Real-Time Monitoring")
tb = add_textbox(s, 1, 1.5, 11.3, 1.5)
tf = tb.text_frame; tf.word_wrap = True
code = ("$ python realtime_focus_monitor.py \\\n"
        "    --data-dir /data/bok/tonight/ \\\n"
        "    --bias-nums 1-10 --dark-nums 21-22 \\\n"
        "    --sci-nums 101-999 --filter all \\\n"
        "    --name-contains OBJECT --incremental")
set_text(tf, code, size=16, color=ORANGE)

add_bullets(s, [
    "Watchdog monitors the data directory for new .fits files",
    "Each new exposure triggers an incremental pipeline run",
    "Cached master calibrations → only ~10 s per new frame",
    "Dashboard updates automatically in the browser",
    "Remote access: HTTP server + ngrok tunnel → share URL with collaborators",
], top=4.0, size=20)

# =====================================================================
# SLIDE 9 — Focus Curve Fitting
# =====================================================================
s = section_slide("Focus Curve Fitting")
add_bullets(s, [
    "Parabola fit in vertex form:  FWHM = A (x − h)² + k",
    "h = best focus position (LVDTC units)",
    "k = minimum FWHM at best focus",
    "σ_h = uncertainty from covariance matrix",
    "R² goodness-of-fit + residual diagnostics",
], width=6.5)

img_path2 = os.path.join(os.path.dirname(__file__), "focus_output", "focus_fit.png")
if os.path.exists(img_path2):
    s.shapes.add_picture(img_path2, Inches(7.5), Inches(1.4), Inches(5.0))
else:
    add_image_placeholder(s, "Insert:\nfocus_fit.png", 7.5, 1.4, 5.0, 4.5)

# =====================================================================
# SLIDE 10 — Focal-Plane Tilt + Focus Solver
# =====================================================================
s = section_slide("Focal-Plane Tilt + Focus Solver")

# ---- Model box (left column) ----------------------------------------
model_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.3), Inches(5.6), Inches(4.4))
model_box.fill.solid(); model_box.fill.fore_color.rgb = DARKBOX
model_box.line.color.rgb = BLUE; model_box.line.width = Pt(1)

tb_model = add_textbox(s, 0.7, 1.45, 5.2, 4.1)
tf = tb_model.text_frame; tf.word_wrap = True
set_text(tf, "Model", size=20, color=BLUE, bold=True)
lines = [
    ("δz(x,y)  =  z₀  +  a·x  +  b·y", ORANGE, True, 18),
    ("", WHITE, False, 10),
    ("FWHM²(x,y)  =  FWHM₀²  +  α · δz²", ORANGE, True, 18),
    ("", WHITE, False, 8),
    ("z₀   overall focus piston (whole camera shift)", WHITE, False, 16),
    ("a, b  tip/tilt gradients along x and y", WHITE, False, 16),
    ("α    defocus sensitivity (FWHM/focus sharpness)", WHITE, False, 16),
    ("FWHM₀  seeing floor (atmosphere + optics)", WHITE, False, 16),
]
for text, col, bold, sz in lines:
    p = tf.add_paragraph(); p.space_before = Pt(4)
    run = p.add_run(); run.text = text
    run.font.size = Pt(sz); run.font.color.rgb = col; run.font.bold = bold

# ---- Procedure box (right column) ----------------------------------
proc_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.4), Inches(1.3), Inches(6.5), Inches(4.4))
proc_box.fill.solid(); proc_box.fill.fore_color.rgb = DARKBOX
proc_box.line.color.rgb = GREY; proc_box.line.width = Pt(1)

tb_proc = add_textbox(s, 6.6, 1.45, 6.1, 4.1)
tf2 = tb_proc.text_frame; tf2.word_wrap = True
set_text(tf2, "Procedure", size=20, color=BLUE, bold=True)
steps = [
    "1  GMM selects stars in each of the 8 amps",
    "2  Median FWHM computed per amp",
    "3  Least-squares fit for 5 params:",
    "     FWHM₀², α, z₀, a, b",
    "4  Corrections at each actuator (A, B, C):",
    "     Δ_act = −(z₀ + a·x_act + b·y_act)",
    "5  Output: tilt_map.png + terminal report",
]
for i, txt in enumerate(steps):
    p = tf2.add_paragraph(); p.space_before = Pt(6 if i == 0 else 8)
    run = p.add_run(); run.text = txt
    run.font.size = Pt(17)
    run.font.color.rgb = ORANGE if txt.startswith(("4 ", "5 ")) else WHITE
    run.font.bold = txt.startswith(("1 ", "2 ", "3 ", "4 ", "5 "))

# ---- Footer --------------------------------------------------------
tb_foot = add_textbox(s, 0.5, 5.9, 12.3, 0.6)
set_text(tb_foot.text_frame,
         "Needs ≥ 5 amps  |  Requires LVDTA / LVDTB / LVDTC header keywords  |  Flag: --solve-tilt",
         size=16, color=GREY, alignment=PP_ALIGN.CENTER)

# ---- Tilt map image (if already generated) -------------------------
tilt_img = os.path.join(os.path.dirname(__file__), "focus_output", "tilt_map.png")
if os.path.exists(tilt_img):
    s.shapes.add_picture(tilt_img, Inches(6.4), Inches(1.3), Inches(6.5), Inches(4.4))

# =====================================================================
# SLIDE 11 (was 10) — Performance
# =====================================================================
s = section_slide("Performance")

perf_data = [
    ("Mode", "121 frames", "Per frame"),
    ("Full run (cold start)", "~25 min", "~12 s"),
    ("Incremental (cached cals)", "~3.5 min", "~10 s new / instant cached"),
    ("Fully cached rerun", "~3.5 min", "PSF thumbnails dominate"),
]
rows, cols = len(perf_data), 3
tbl = s.shapes.add_table(rows, cols, Inches(1.5), Inches(1.6), Inches(10.3), Inches(2.5)).table
for c in range(cols):
    tbl.columns[c].width = Inches([4.0, 3.0, 3.3][c])
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = perf_data[r][c]
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARKBOX if r == 0 else BG
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(20)
                run.font.color.rgb = BLUE if r == 0 else WHITE
                run.font.bold = (r == 0)

add_bullets(s, [
    "Cache: master bias/flat per band per amp (.npy), per-file catalogs (.fits) + cutouts (.npz)",
    "Tested on MacBook — at the telescope (local data) → even faster",
], top=4.6, size=20)

# =====================================================================
# SLIDE 11 — CLI Usage
# =====================================================================
s = section_slide("Usage")
tb = add_textbox(s, 0.8, 1.5, 11.7, 5.5)
tf = tb.text_frame; tf.word_wrap = True
code1 = ("# Single-band run\n"
         "$ python focus_pipeline.py \\\n"
         "    --data-dir /data/tonight/ \\\n"
         "    --filter r \\\n"
         "    --bias-nums 1-10 --dark-nums 21-22 \\\n"
         "    --flat-nums 91-100 --sci-nums 101-200 \\\n"
         "    --incremental\n\n"
         "# Multi-band: auto-discover all bands + flats\n"
         "$ python focus_pipeline.py \\\n"
         "    --data-dir /data/tonight/ \\\n"
         "    --filter all \\\n"
         "    --bias-nums 1-10 --dark-nums 21-22 \\\n"
         "    --sci-nums 101-300 \\\n"
         "    --incremental")
set_text(tf, code1, size=16, color=ORANGE)

tb2 = add_textbox(s, 1, 6.2, 11.3, 0.6)
set_text(tb2.text_frame,
         "Key flags: --focus-key  --time-key  --date-key  --airmass-key  --pixscale  --threshold  --amps",
         size=16, color=GREY, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 12 — Science Insights
# =====================================================================
s = section_slide("Science Insights from Nov 11, 2025 Data")
add_bullets(s, [
    "Seeing ranged 1.0\" – 2.1\" (r-band) over ~7 hours",
    "g-band FWHM consistently ~19% larger than r-band at same time:",
    "     ~5% from atmospheric chromaticity (Fried's law: FWHM ∝ λ^(−1/5))",
    "     ~14% from chromatic defocus — refractive corrector optimized for r",
    "Ellipticity stable at ~0.1 – 0.15 → good tracking all night",
    "Airmass 1.02 – 1.26 → moderate zenith angle range",
])

# =====================================================================
# SLIDE 13 — Future Directions
# =====================================================================
s = section_slide("Future Directions")
add_bullets(s, [
    "🔮  Seeing prediction: time-series forecasting of FWHM trends",
    "🎯  Auto-focus closed loop: trigger refocus when FWHM > threshold",
    "🌡️  Temperature correlation: model focus drift vs ambient temperature",
    "📱  Alerts: Slack / email notification when seeing degrades",
    "🔭  Generalize to other telescopes (MMT, Kuiper, LBT) — only header keyword config needed",
    "📦  Package: pip-installable with proper config files",
])

# =====================================================================
# SLIDE 14 — Summary
# =====================================================================
s = section_slide("Summary")
add_bullets(s, [
    "✅  Automated end-to-end: raw FITS → interactive dashboard",
    "✅  Real-time: ~10 s per new frame, 30 s auto-refresh",
    "✅  Multi-band: auto-discovers bands and flats",
    "✅  Interactive: hover for PSF contours, zoom, pan",
    "✅  Remote access: HTTP server + ngrok",
    "✅  Quantitative: FWHM in arcsec, focus curve fit, GMM star selection",
])
tb = add_textbox(s, 1, 5.5, 11.3, 0.6)
set_text(tb.text_frame, "github.com/jenny/autofocus", size=28, color=ORANGE, bold=True,
         alignment=PP_ALIGN.CENTER)
tb2 = add_textbox(s, 1, 6.3, 11.3, 0.5)
set_text(tb2.text_frame, "Python 3.8+  |  astropy · SEP · plotly · scikit-learn · watchdog",
         size=16, color=GREY, alignment=PP_ALIGN.CENTER)

# =====================================================================
# Save — ensure file is NOT read-only
# =====================================================================
outpath = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AutoFocus_Slides.pptx")

# Remove existing file first (in case it has read-only OS permissions)
if os.path.exists(outpath):
    os.chmod(outpath, 0o644)
    os.remove(outpath)

# Clear any PowerPoint security/read-only flags in the core properties
prs.core_properties.revision = 1
prs.core_properties.last_modified_by = ""

prs.save(outpath)

# Ensure the file has full read-write permissions
os.chmod(outpath, 0o644)
print(f"Saved (editable) → {outpath}")
