#!/usr/bin/env python3
"""Generate tilt-solver / actuator-correction explanation slides as tilt_map_slides.pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------- colour palette ----------
BG      = RGBColor(0x0D, 0x11, 0x17)
BLUE    = RGBColor(0x58, 0xA6, 0xFF)
ORANGE  = RGBColor(0xFF, 0x8C, 0x00)
WHITE   = RGBColor(0xC9, 0xD1, 0xD9)
GREY    = RGBColor(0x8B, 0x94, 0x9E)
DARKBOX = RGBColor(0x16, 0x1B, 0x22)
GREEN   = RGBColor(0x3F, 0xB9, 0x50)
RED     = RGBColor(0xF8, 0x51, 0x49)
YELLOW  = RGBColor(0xE3, 0xB3, 0x41)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ---- helpers ----
def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))


def set_text(tf, text, size=24, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_para(tf, text, size=22, color=WHITE, bold=False,
             space_before=6, alignment=None):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    if alignment is not None:
        p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    tb = add_textbox(slide, 1, 2.0, 11.3, 1.5)
    set_text(tb.text_frame, title, size=44, color=BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)
    if subtitle:
        tb2 = add_textbox(slide, 1, 3.7, 11.3, 1.0)
        set_text(tb2.text_frame, subtitle, size=24, color=GREY,
                 alignment=PP_ALIGN.CENTER)
    return slide


def section_slide(heading):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    tb = add_textbox(slide, 0.8, 0.4, 11.7, 0.9)
    set_text(tb.text_frame, heading, size=36, color=BLUE, bold=True)
    return slide


def add_bullets(slide, items, left=1.0, top=1.6, width=11.3, height=5.0,
                size=22, color=WHITE):
    tb = add_textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            set_text(tf, "•  " + item, size=size, color=color)
        else:
            add_para(tf, "•  " + item, size=size, color=color, space_before=10)
    return tb


def add_dark_box(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARKBOX
    shape.line.color.rgb = GREY
    shape.line.width = Pt(1)
    return shape


def add_labeled_box(slide, label, value, left, top, width=3.2, height=1.3,
                    val_size=32, label_size=16):
    shape = add_dark_box(slide, left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    set_text(tf, value, size=val_size, color=ORANGE, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_para(tf, label, size=label_size, color=GREY, alignment=PP_ALIGN.CENTER)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    return shape


# =====================================================================
# SLIDE 1 — Title
# =====================================================================
s = title_slide(
    "Tilt Map & Actuator Corrections",
    "How Δ_A, Δ_B, Δ_C are computed from the FWHM map")
tb = add_textbox(s, 1, 4.9, 11.3, 0.8)
set_text(tb.text_frame,
         "Bok 90Prime  |  solve_tilt_focus  |  focus_pipeline.py",
         size=18, color=GREY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 2 — The Defocus Plane Model
# =====================================================================
s = section_slide("Step 1 — Fit a Defocus Plane to the FWHM Map")

# Model equation box
eq_box = add_dark_box(s, 0.8, 1.3, 11.7, 1.8)
tf = eq_box.text_frame
tf.word_wrap = True
set_text(tf,
         "FWHM²(x,y)  =  FWHM₀²  +  α · δz(x,y)²",
         size=26, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_para(tf,
         "δz(x,y)  =  z₀  +  a·x  +  b·y",
         size=24, color=WHITE, alignment=PP_ALIGN.CENTER, space_before=8)

# Parameter table
params = [
    ("FWHM₀", "seeing floor — set by atmosphere, same on all amps"),
    ("α",     "defocus sensitivity (FWHM² growth per defocus unit²)"),
    ("z₀",    "piston — overall focus offset at the field centre"),
    ("a, b",  "tip / tilt slopes across the focal plane"),
]
tb2 = add_textbox(s, 1.0, 3.3, 11.3, 3.5)
tf2 = tb2.text_frame
tf2.word_wrap = True
set_text(tf2, "Free parameters fitted by least-squares to 8 per-amp median FWHM values:",
         size=19, color=GREY)
for sym, desc in params:
    add_para(tf2, f"   {sym}  —  {desc}", size=20, color=WHITE, space_before=9)


# =====================================================================
# SLIDE 3 — Evaluate plane at actuator positions
# =====================================================================
s = section_slide("Step 2 — Evaluate the Plane at Each Actuator")

# Diagram: three actuator positions
act_data = [
    ("A",  2.5, 4.4, BLUE),
    ("B",  6.2, 2.4, GREEN),
    ("C", 10.0, 4.4, ORANGE),
]
for name, lx, ly, col in act_data:
    circ = s.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(lx - 0.45), Inches(ly - 0.45), Inches(0.9), Inches(0.9))
    circ.fill.solid()
    circ.fill.fore_color.rgb = col
    circ.line.fill.background()
    tf = circ.text_frame
    set_text(tf, name, size=22, color=BG, bold=True, alignment=PP_ALIGN.CENTER)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    lbl = add_textbox(s, lx - 1.1, ly + 0.52, 2.2, 0.5)
    set_text(lbl.text_frame,
             f"(x_{name}, y_{name})",
             size=16, color=col, alignment=PP_ALIGN.CENTER)
    eq = add_textbox(s, lx - 1.4, ly + 1.0, 2.8, 0.55)
    set_text(eq.text_frame,
             f"δz_{name} = z₀ + a·x_{name} + b·y_{name}",
             size=15, color=GREY, alignment=PP_ALIGN.CENTER)

# Arrow lines between actuators (text stand-ins)
arr = add_textbox(s, 3.1, 4.05, 3.0, 0.5)
set_text(arr.text_frame, "120°  ←→  120°", size=16, color=GREY, alignment=PP_ALIGN.CENTER)

# Caption
tb = add_textbox(s, 1.0, 1.5, 11.3, 0.7)
set_text(tb.text_frame,
         "Actuators A, B, C are 120° apart on a circle (_actuator_xy).  "
         "δz_k is just the plane value at that (x, y).",
         size=20, color=WHITE)


# =====================================================================
# SLIDE 4 — The Correction Formula
# =====================================================================
s = section_slide("Step 3 — Correction = Move That Cancels the Local Defocus")

# Big formula
eq_box = add_dark_box(s, 1.2, 1.3, 10.9, 2.0)
tf = eq_box.text_frame
tf.word_wrap = True
set_text(tf,
         "Δ_k  =  −δz_k  =  −(z₀ + a·x_k + b·y_k)",
         size=30, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_para(tf,
         "LVDT_k(opt)  =  LVDT_k(current)  +  Δ_k",
         size=24, color=WHITE, alignment=PP_ALIGN.CENTER, space_before=10)

# Code snippet box
code_box = add_dark_box(s, 1.2, 3.55, 10.9, 2.5)
tf2 = code_box.text_frame
tf2.word_wrap = True
set_text(tf2,
         "# focus_pipeline.py  (solve_tilt_focus)",
         size=15, color=GREY)
for line in [
    "for name, (ax, ay) in act_xy.items():",
    "    defocus_at_act       = z0 + a * ax + b * ay",
    "    corrections[name]   = -defocus_at_act",
    "    optimal_lvdt[name]  = current_lvdt[name] - defocus_at_act",
]:
    add_para(tf2, line, size=16, color=GREEN, space_before=4)


# =====================================================================
# SLIDE 5 — Geometry Intuition
# =====================================================================
s = section_slide("Geometry Intuition")

cases = [
    (BLUE,   "Pure Piston",
     "a = b = 0,  z₀ ≠ 0",
     ["All three Δ's are equal",
      "Push all three actuators by the same amount",
      "→ uniform focus shift, no tilt change"]),
    (ORANGE, "Pure Tilt (along x)",
     "a ≠ 0,  z₀ = b = 0",
     ["+x actuator gets a negative Δ",
      "−x actuator gets a positive Δ",
      "→ mirror tips to flatten the FWHM gradient"]),
    (GREEN,  "General Case",
     "z₀ ≠ 0,  a ≠ 0,  b ≠ 0",
     ["Linear combination of piston + tip + tilt",
      "Each actuator gets a different correction",
      "→ simultaneous focus + tilt correction"]),
]
for i, (col, title, subtitle, bullets) in enumerate(cases):
    bx = 0.4 + i * 4.3
    box = add_dark_box(s, bx, 1.3, 4.05, 5.6)
    tf = box.text_frame
    tf.word_wrap = True
    set_text(tf, title, size=22, color=col, bold=True, alignment=PP_ALIGN.CENTER)
    add_para(tf, subtitle, size=17, color=YELLOW,
             alignment=PP_ALIGN.CENTER, space_before=4)
    for b in bullets:
        add_para(tf, "•  " + b, size=16, color=WHITE, space_before=8)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER


# =====================================================================
# SLIDE 6 — Caveats: Sign Degeneracy
# =====================================================================
s = section_slide("Caveat 1 — Sign Degeneracy")

# Main point
warn_box = add_dark_box(s, 0.8, 1.3, 11.7, 1.5)
tf = warn_box.text_frame
tf.word_wrap = True
set_text(tf,
         "Only δz² enters the model  →  (z₀, a, b) and −(z₀, a, b) give identical fits",
         size=22, color=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)
add_para(tf,
         "The solver tries both seeds (+ and −) and keeps the lower-residual one",
         size=19, color=WHITE, alignment=PP_ALIGN.CENTER, space_before=6)

add_bullets(s, [
    "Per-frame sign flips can still happen if the residuals are nearly equal",
    "A sign flip inverts all three corrections:  Δ_A → −Δ_A etc.",
    "Individual frames can silently point the wrong way",
], top=3.0, size=21)

tb = add_textbox(s, 0.8, 5.5, 11.7, 0.75)
set_text(tb.text_frame,
         "Mitigation:  average corrections across ≥5 high-R² frames  "
         "— sign-consistent frames will dominate",
         size=20, color=GREEN, bold=True)


# =====================================================================
# SLIDE 7 — Caveats: Seeing ↔ Piston Degeneracy
# =====================================================================
s = section_slide("Caveat 2 — Seeing ↔ Piston Degeneracy (single frame)")

# Explanation
tb = add_textbox(s, 0.8, 1.4, 11.7, 1.0)
set_text(tb.text_frame,
         "Atmosphere (FWHM₀) affects all amps equally — indistinguishable from uniform piston in one frame.",
         size=21, color=WHITE)

# 4-row reliability table
headers = ["Quantity",                        "Single-frame reliability"]
rows = [
    ("Tilt slopes  a, b",                     "✅  Robust — from FWHM differences between amps"),
    ("Tilt-only deltas  Δ_k − mean(Δ)",       "✅  Robust"),
    ("Piston z₀  (mean of Δ_A, Δ_B, Δ_C)",   "⚠️  Degenerate with seeing — unreliable"),
    ("Best-focus LVDT from tilt fit alone",   "❌  Do NOT trust"),
]
col_w = [4.4, 7.0]
col_x = [0.8, 5.4]
row_h = 0.72
row_y_start = 2.55

# header row
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    b = add_dark_box(s, cx, row_y_start, cw, row_h)
    b.line.color.rgb = BLUE
    tf = b.text_frame
    set_text(tf, hdr, size=18, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER

for ri, (label, note) in enumerate(rows):
    ry = row_y_start + (ri + 1) * row_h
    col_colors = [WHITE, GREEN if "✅" in note else (YELLOW if "⚠" in note else RED)]
    for ci, (cx, cw, text, col) in enumerate(
            zip(col_x, col_w, [label, note], col_colors)):
        b = add_dark_box(s, cx, ry, cw, row_h)
        tf = b.text_frame
        set_text(tf, text, size=16, color=col, alignment=PP_ALIGN.CENTER)
        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.CENTER

tb2 = add_textbox(s, 0.8, 6.75, 11.7, 0.55)
set_text(tb2.text_frame,
         "For overall focus: use the LVDT-vs-FWHM parabola fit, not the tilt solver.",
         size=18, color=ORANGE, bold=True)


# =====================================================================
# SLIDE 8 — Global Fit: Breaking the Degeneracy
# =====================================================================
s = section_slide("--global-tilt-fit  —  Breaking the Seeing/Piston Degeneracy")

# Command box
cmd_box = add_dark_box(s, 0.8, 1.3, 11.7, 0.85)
tf = cmd_box.text_frame
set_text(tf,
         "python focus_pipeline.py  [...]  --solve-tilt  --global-tilt-fit",
         size=20, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
for p in tf.paragraphs:
    p.alignment = PP_ALIGN.CENTER

# Explanation
add_bullets(s, [
    "Fits one (FWHM₀, α) jointly across all frames in the focus scan",
    "Per-frame (z₀, a, b) can vary freely — but seeing must be consistent",
    "Because defocus changes between LVDT steps, the degeneracy is broken",
    "Writes  tilt_global.json  (global summary) + per-frame tilt_result_<num>.json",
    "Requires ≥ 5 frames at varying LVDT positions",
], top=2.4, size=20)

# Stat boxes
for i, (val, lbl) in enumerate([
        ("FWHM₀, α", "shared\nacross frames"),
        ("z₀, a, b", "per-frame\n(free)"),
        ("≥ 5 frames", "recommended\nminimum")]):
    add_labeled_box(s, lbl, val, 1.0 + i * 4.1, 5.8, width=3.6, height=1.4,
                    val_size=22, label_size=15)


# =====================================================================
# SLIDE 9 — Recommended Workflow
# =====================================================================
s = section_slide("Recommended Workflow")

steps = [
    ("1", "Focus scan\n~10 exposures\nstepped in LVDT",           BLUE),
    ("2", "--solve-tilt\n--global-tilt-fit",                       ORANGE),
    ("3", "Piston\nfrom LVDT-vs-FWHM\nparabola",                   GREEN),
    ("4", "Tilt-only deltas\nΔ_k − mean(Δ)\nfrom tilt_result",    YELLOW),
    ("5", "Command\nactuators\npiston + tilt",                      WHITE),
]
for i, (num, label, col) in enumerate(steps):
    bx = 0.5 + i * 2.5
    box = add_dark_box(s, bx, 1.5, 2.25, 3.2)
    box.line.color.rgb = col
    tf = box.text_frame
    tf.word_wrap = True
    set_text(tf, num, size=36, color=col, bold=True, alignment=PP_ALIGN.CENTER)
    add_para(tf, label, size=16, color=WHITE, alignment=PP_ALIGN.CENTER, space_before=8)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    if i < len(steps) - 1:
        arr = add_textbox(s, bx + 2.25, 2.7, 0.3, 0.6)
        set_text(arr.text_frame, "→", size=28, color=GREY, alignment=PP_ALIGN.CENTER)

# Code snippet for tilt-only delta
code_box = add_dark_box(s, 0.8, 5.0, 11.7, 2.0)
tf2 = code_box.text_frame
tf2.word_wrap = True
set_text(tf2, "# Tilt-only deltas (removes piston ambiguity)", size=15, color=GREY)
for line in [
    "d    = tilt_result['corrections']          # {'A': ..., 'B': ..., 'C': ...}",
    "mean = (d['A'] + d['B'] + d['C']) / 3",
    "tilt_only = {k: d[k] - mean for k in d}   # zero-mean actuator deltas",
]:
    add_para(tf2, line, size=16, color=GREEN, space_before=5)


# =====================================================================
# SLIDE 10 — Summary
# =====================================================================
s = section_slide("Summary")

add_bullets(s, [
    "FWHM²(x,y) = FWHM₀² + α·δz²  fitted to 8-amp map → piston z₀, tilt (a, b)",
    "Δ_k = −(z₀ + a·x_k + b·y_k)  — correction per actuator, LVDT units",
    "Tilt slopes (a, b) are reliable from a single frame; piston z₀ is not",
    "Use --global-tilt-fit on a focus scan (≥5 frames) to break seeing/piston degeneracy",
    "Final actuator move = piston (from focus curve) + tilt-only deltas (from tilt solver)",
], size=22)

tb = add_textbox(s, 1, 6.2, 11.3, 0.6)
set_text(tb.text_frame,
         "solve_tilt_focus  &  solve_tilt_focus_global  |  focus_pipeline.py",
         size=18, color=GREY, alignment=PP_ALIGN.CENTER)


# =====================================================================
# Save
# =====================================================================
outpath = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "tilt_map_slides.pptx")
if os.path.exists(outpath):
    os.chmod(outpath, 0o644)
    os.remove(outpath)

prs.core_properties.revision = 1
prs.core_properties.last_modified_by = ""

prs.save(outpath)
os.chmod(outpath, 0o644)
print(f"Saved → {outpath}")
